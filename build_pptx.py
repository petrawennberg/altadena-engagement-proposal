#!/usr/bin/env python3
"""
Build a .pptx version of the CAPS deck for Google Slides import.

Visual design follows the joinalting brand:
- Cream background, forest ink text
- Coral + sage accent colors
- DM Serif Display (headings), DM Sans (body)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============ BRAND TOKENS ============
CREAM = RGBColor(0xFA, 0xF6, 0xF0)
CREAM_DEEP = RGBColor(0xF0, 0xEA, 0xE0)
FOREST = RGBColor(0x2E, 0x3D, 0x26)
FOREST_SOFT = RGBColor(0x5A, 0x69, 0x53)
CORAL = RGBColor(0xE5, 0x90, 0x5A)
CORAL_TINT = RGBColor(0xFB, 0xEB, 0xDC)
SAGE = RGBColor(0x7A, 0x8A, 0x6E)
SAGE_TINT = RGBColor(0xE8, 0xED, 0xE3)
PINK = RGBColor(0xE8, 0x5A, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "DM Serif Display"
SANS = "DM Sans"

# ============ HELPERS ============
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, *,
                  font=SANS, size=14, color=FOREST, bold=False, italic=False,
                  align="left", vertical="top", spacing=1.2, letter_spacing=0,
                  uppercase=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    if vertical == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif vertical == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb

def add_rect(slide, left, top, width, height, fill_color, *, line_color=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if line_color:
        rect.line.color.rgb = line_color
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect

def add_rounded_rect(slide, left, top, width, height, fill_color, *, line_color=None, line_width=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.adjustments[0] = 0.06
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill_color
    if line_color:
        rect.line.color.rgb = line_color
        if line_width:
            rect.line.width = line_width
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect

def add_slide_tag(slide, text, top=Inches(0.55)):
    add_text_box(slide, Inches(0.7), top, Inches(8), Inches(0.3),
                 text, font=SANS, size=11, color=CORAL, bold=True,
                 uppercase=True, spacing=1.0)

def add_slide_title(slide, text, top=Inches(0.85), size=32):
    add_text_box(slide, Inches(0.7), top, Inches(12), Inches(1.2),
                 text, font=SERIF, size=size, color=FOREST, spacing=1.1)

def add_footer(slide):
    # Thin sage line above footer
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(7.05), Inches(12.63), Inches(7.05))
    line.line.color.rgb = SAGE_TINT
    line.line.width = Pt(0.5)
    add_text_box(slide, Inches(0.7), Inches(7.13), Inches(12), Inches(0.3),
                 "ALTADENA VILLAGE  ·  COLAB  ·  ALTING  ·  ALTAGETHER",
                 font=SANS, size=9, color=FOREST_SOFT, bold=False,
                 align="center", spacing=1.0)

# ============ BUILD DECK ============
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

blank_layout = prs.slide_layouts[6]

# ----- SLIDE 1: TITLE -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
# Coral accent
add_rect(s, Inches(6.16), Inches(2.5), Inches(1), Inches(0.05), CORAL)
# Eyebrow
add_text_box(s, Inches(0.7), Inches(2.1), Inches(11.93), Inches(0.3),
             "A COLLABORATION", font=SANS, size=11, color=CORAL, bold=True,
             align="center", letter_spacing=6, uppercase=True)
# Main title
add_text_box(s, Inches(1), Inches(2.8), Inches(11.33), Inches(2.5),
             "A community engagement strategy for Altadena rebuild",
             font=SERIF, size=48, color=FOREST, align="center", spacing=1.05)
# Partners
add_text_box(s, Inches(0.7), Inches(5.4), Inches(11.93), Inches(0.4),
             "Altadena Village  ·  CoLab  ·  Alting  ·  Altagether",
             font=SANS, size=14, color=FOREST, bold=True, align="center")
# Meta
add_text_box(s, Inches(0.7), Inches(6.3), Inches(11.93), Inches(0.3),
             "PRESENTED TO CAPS  ·  MAY 21, 2026",
             font=SANS, size=10, color=FOREST_SOFT, align="center",
             uppercase=True)

# ----- SLIDE 2: WHY (headline) -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "Why we're here")
add_text_box(s, Inches(0.7), Inches(2), Inches(11.93), Inches(3),
             "Altadena's commercial and public spaces are being rebuilt. There's no process for the community to shape what they become and see the impact of their input.",
             font=SERIF, size=36, color=FOREST, spacing=1.15)
add_footer(s)

# ----- SLIDE 3: TIMELINE (Gantt) -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "Why we're here")
add_slide_title(s, "Separate schedules, no coordination", size=28)

# Gantt area
chart_left = Inches(0.7)
chart_top = Inches(1.8)
chart_width = Inches(12.63)
chart_height = Inches(4.8)
label_col_width = Inches(2.0)
months_area_left = chart_left + label_col_width
months_area_width = chart_width - label_col_width
num_months = 22
month_width = months_area_width / num_months

# Year headers
year_2026_width = month_width * 10
year_2027_width = month_width * 12
year_h = Inches(0.25)
yr1 = add_rect(s, months_area_left, chart_top, year_2026_width - Inches(0.02), year_h, SAGE_TINT)
add_text_box(s, months_area_left, chart_top, year_2026_width, year_h, "2026",
             font=SANS, size=9, color=FOREST, bold=True, align="center", vertical="middle", uppercase=True)
yr2 = add_rect(s, months_area_left + year_2026_width, chart_top, year_2027_width, year_h, SAGE_TINT)
add_text_box(s, months_area_left + year_2026_width, chart_top, year_2027_width, year_h, "2027",
             font=SANS, size=9, color=FOREST, bold=True, align="center", vertical="middle", uppercase=True)

# Month labels
months = ["Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec",
          "Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
month_y = chart_top + year_h + Inches(0.05)
for i, m in enumerate(months):
    add_text_box(s, months_area_left + i*month_width, month_y, month_width, Inches(0.22), m,
                 font=SANS, size=8, color=FOREST_SOFT, bold=True, align="center", uppercase=True)

# Underline
under_y = month_y + Inches(0.24)
under = slide_line = s.shapes.add_connector(1, months_area_left, under_y, months_area_left + months_area_width, under_y)
under.line.color.rgb = SAGE_TINT
under.line.width = Pt(0.5)

# Rows
rows_top = under_y + Inches(0.08)
row_h = Inches(0.34)
row_gap = Inches(0.06)

gantt_rows = [
    ("CIP", "LA County Public Works", [
        (1, 3, "Survey", CORAL, "solid"),
        (4, 15, "Comment period TBD", CORAL, "tbd"),
    ]),
    ("Regional Planning", "WSGVAP refresh", [
        (5, 20, "Refresh cycle TBD", SAGE, "tbd-sage"),
    ]),
    ("Sanitation", "LA County sewer", [
        (1, 23, "Ongoing rebuild interaction", SAGE, "solid"),
    ]),
    ("Water utilities", "Las Flores · Lincoln · Rubio Cañon", [
        (1, 23, "Ongoing restoration & planning", SAGE, "solid"),
    ]),
    ("Flood Control", "LA County stormwater", [
        (4, 17, "Stormwater windows TBD", SAGE, "tbd-sage"),
    ]),
    ("Edison / SCE", "", [
        (1, 23, "Undergrounding & restoration — ongoing", SAGE, "solid"),
    ]),
    ("SoCalGas", "", [
        (1, 23, "Gas restoration ongoing", SAGE, "solid"),
    ]),
    ("PUSD", "School-site rebuild", [
        (5, 18, "Decisions TBD", CORAL, "tbd"),
    ]),
    ("Parks & Rec", "LA County multi-benefit", [
        (6, 20, "Parks & green corridors TBD", SAGE, "tbd-sage"),
    ]),
    ("Tree restoration", "LA County · TreePeople · Altadena Bloom", [
        (6, 23, "Urban canopy restoration TBD", SAGE, "tbd-sage"),
    ]),
    ("Commercial developers", "Various private projects", [
        (1, 12, "Active rebuild projects ongoing", SAGE, "solid"),
        (12, 23, "New permits & projects TBD", SAGE, "tbd-sage"),
    ]),
    ("Paradigm", "Edison's outreach consultant", [
        (1, 23, "Engagement windows TBD", CORAL, "tbd"),
    ]),
    ("Toole", "LA County's planning consultant", [
        (1, 6, "Master plan in progress", FOREST, "solid"),
        (6, 23, "No engagement scheduled", CORAL, "tbd"),
    ]),
]

# Note: 13 rows. Fit into chart_height = 4.8" total ~13*(row_h+gap) = 13*(0.34+0.06) = 5.2"
# Need to shrink rows. Let's use row_h=0.28, gap=0.04 = 13*0.32 = 4.16". Fits.
row_h = Inches(0.28)
row_gap = Inches(0.04)

for ri, (label, sublabel, bars) in enumerate(gantt_rows):
    y = rows_top + ri*(row_h + row_gap)
    # Label
    add_text_box(s, chart_left, y, label_col_width - Inches(0.1), Inches(0.18), label,
                 font=SANS, size=10, color=FOREST, bold=True, vertical="middle")
    if sublabel:
        add_text_box(s, chart_left, y + Inches(0.16), label_col_width - Inches(0.1), Inches(0.15), sublabel,
                     font=SANS, size=7, color=FOREST_SOFT, vertical="middle")
    # Track background
    track = add_rounded_rect(s, months_area_left, y, months_area_width, row_h,
                              RGBColor(0xEE, 0xF0, 0xEA))
    # Bars
    for (col_start, col_end, text, color, style) in bars:
        bar_left = months_area_left + (col_start - 1) * month_width + Inches(0.01)
        bar_right = months_area_left + (col_end - 1) * month_width - Inches(0.01)
        bar_w = bar_right - bar_left
        if style == "solid":
            bar = add_rounded_rect(s, bar_left, y + Inches(0.02), bar_w, row_h - Inches(0.04), color)
            add_text_box(s, bar_left + Inches(0.06), y + Inches(0.02), bar_w - Inches(0.12), row_h - Inches(0.04),
                         text, font=SANS, size=8, color=WHITE, bold=True, vertical="middle")
        elif style in ("tbd", "tbd-sage"):
            border_color = color
            bar = add_rounded_rect(s, bar_left, y + Inches(0.02), bar_w, row_h - Inches(0.04), CREAM,
                                    line_color=border_color, line_width=Pt(1.5))
            add_text_box(s, bar_left + Inches(0.06), y + Inches(0.02), bar_w - Inches(0.12), row_h - Inches(0.04),
                         text, font=SANS, size=8, color=color, bold=True, vertical="middle", italic=True)

# Punch line about PFA
pfa_y = rows_top + 13*(row_h + row_gap) + Inches(0.08)
pfa_box = add_rounded_rect(s, chart_left, pfa_y, chart_width, Inches(0.55), CORAL_TINT)
add_text_box(s, chart_left + Inches(0.2), pfa_y + Inches(0.04), Inches(1.4), Inches(0.18),
             "AND THE PFA", font=SANS, size=8, color=CORAL, bold=True, uppercase=True)
add_text_box(s, chart_left + Inches(0.2), pfa_y + Inches(0.22), chart_width - Inches(0.4), Inches(0.32),
             "The Public Financing Authority has convened — but hasn't held the meetings required by its own bylaws. The engagement standard adopted at those meetings will be the default for every project above.",
             font=SANS, size=9, color=FOREST, spacing=1.3)

add_footer(s)

# ----- SLIDE 4: WHAT'S MISSING (5 values) -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "What's missing")
add_slide_title(s, "Five gaps — and the values that close them", size=28)

values = [
    ("01", "Local insight", "Residents carry context the County doesn't have, with no channel for it to land."),
    ("02", "Local agency", "Commercial owners have coordinated a shared vision, but no defined path to the County's infrastructure rollout."),
    ("03", "Coordination", "Agencies plan in parallel, on separate schedules, with no shared way to coordinate."),
    ("04", "Accountability", "No way to track whether community input shaped what actually gets built."),
    ("05", "Trauma-informed", "No engagement design built for a community that's been through this much."),
]
row_top = Inches(2.1)
row_h = Inches(0.92)
for i, (num, name, desc) in enumerate(values):
    y = row_top + i*row_h
    # Number
    add_text_box(s, Inches(0.7), y, Inches(0.8), Inches(0.6), num,
                 font=SERIF, size=32, color=CORAL, vertical="top")
    # Name
    add_text_box(s, Inches(1.7), y, Inches(3.5), Inches(0.4), name,
                 font=SANS, size=18, color=FOREST, bold=True)
    # Desc
    add_text_box(s, Inches(5.4), y, Inches(7.5), Inches(0.7), desc,
                 font=SANS, size=14, color=FOREST_SOFT, spacing=1.4)
    # Divider line
    if i < len(values) - 1:
        ln = s.shapes.add_connector(1, Inches(0.7), y + row_h - Inches(0.05), Inches(12.63), y + row_h - Inches(0.05))
        ln.line.color.rgb = SAGE_TINT
        ln.line.width = Pt(0.5)
add_footer(s)

# ----- SLIDE 5: PROPOSAL (5-step framework) -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "What we're proposing")
add_slide_title(s, "Five steps — content flows down, consensus flows up", size=24)

# Layout: 2 columns top, full-width bottom
card_top = Inches(1.7)
card_h_top = Inches(3.6)
gap_v = Inches(0.15)
card_bot_top = card_top + card_h_top + gap_v
card_h_bot = Inches(1.45)

col_left = Inches(0.7)
col_w = Inches(6.18)
gap_h = Inches(0.18)
col_right = col_left + col_w + gap_h

arrow_w = Inches(0.4)

# Helper to draw a flow card
def draw_flow_card(slide, x, y, w, h, arrow_bg, arrow_color, arrow_char, label, steps, dark=False):
    # Card background
    card_bg = WHITE if not dark else FOREST
    add_rounded_rect(slide, x, y, w, h, card_bg)
    # Arrow column
    add_rounded_rect(slide, x, y, arrow_w, h, arrow_bg)
    # Arrows
    arrow_text_color = arrow_color
    arrow_y_start = y + Inches(0.5)
    arrow_gap = Inches(0.5)
    for ai in range(2):
        add_text_box(slide, x, arrow_y_start + ai*arrow_gap, arrow_w, Inches(0.3),
                     arrow_char, font=SANS, size=16, color=arrow_text_color, bold=True, align="center")
    # Content area
    content_left = x + arrow_w + Inches(0.2)
    content_w = w - arrow_w - Inches(0.3)
    text_color = FOREST if not dark else CREAM
    soft_color = FOREST_SOFT if not dark else RGBColor(0xC8, 0xCC, 0xC0)
    add_text_box(slide, content_left, y + Inches(0.15), content_w, Inches(0.25),
                 label, font=SANS, size=9, color=arrow_color, bold=True, uppercase=True)
    # Steps
    step_y = y + Inches(0.5)
    for (num, name, desc, who) in steps:
        step_h = (h - Inches(0.5)) / max(len(steps), 1) - Inches(0.05)
        # Number
        add_text_box(slide, content_left, step_y, Inches(0.55), Inches(0.4),
                     num, font=SERIF, size=22, color=arrow_color)
        # Name
        add_text_box(slide, content_left + Inches(0.55), step_y, content_w - Inches(0.55), Inches(0.28),
                     name, font=SANS, size=14, color=text_color, bold=True)
        # Description
        add_text_box(slide, content_left + Inches(0.55), step_y + Inches(0.3), content_w - Inches(0.55), Inches(0.65),
                     desc, font=SANS, size=10, color=text_color, spacing=1.35)
        # Who
        if who:
            add_text_box(slide, content_left + Inches(0.55), step_y + step_h - Inches(0.25), content_w - Inches(0.55), Inches(0.2),
                         "WHO  " + who, font=SANS, size=8, color=soft_color, italic=True)
        step_y += step_h + Inches(0.05)

# Down column (steps 01, 02)
down_steps = [
    ("01", "Definition",
     "Sponsors define their decision based on community standards — problem, ask, geographic area, timeline, and multi-modal delivery (text · video · Spanish · walking tours) with expert access.",
     "County agencies · utilities · developers · consultants"),
    ("02", "Training & Distribution",
     "Through existing channels — Altagether's 150+ block captains and the commercial-owner network, trained to facilitate with on-call backup.",
     "Altagether captains · commercial-owner network · CoLab (training)"),
]
draw_flow_card(s, col_left, card_top, col_w, card_h_top, CORAL_TINT, CORAL, "↓",
                "↓ Content flows down", down_steps)

# Up column (steps 05, 04)
up_steps = [
    ("05", "Shared record",
     "The record goes public — tracked on a dashboard, visible to all stakeholders.",
     "All stakeholders see · working group maintains"),
    ("04", "Aggregation",
     "Block deliberations compile into a coherent record — organized by area, corridor, zone, and by the sponsor's own categories. Private dialogue stays private; outcomes become public.",
     "Captains record · CoLab captures · Alting synthesizes"),
]
draw_flow_card(s, col_right, card_top, col_w, card_h_top, SAGE_TINT, SAGE, "↑",
                "↑ Consensus flows up", up_steps)

# Bottom (Deliberation - step 03)
bot_steps = [
    ("03", "Facilitating Discussion",
     "Block captains facilitate time-bound community deliberation at the block level.",
     "Block captains lead · Altadena public participates · local facilitators support"),
]
# For dark card, arrow column is forest-soft
add_rounded_rect(s, col_left, card_bot_top, col_w*2 + gap_h, card_h_bot, FOREST)
add_rounded_rect(s, col_left, card_bot_top, arrow_w, card_h_bot, FOREST_SOFT)
# Dots in arrow col
for di in range(2):
    add_text_box(s, col_left, card_bot_top + Inches(0.35) + di*Inches(0.35), arrow_w, Inches(0.25),
                 "•", font=SANS, size=18, color=CREAM, bold=True, align="center")
# Content
content_left = col_left + arrow_w + Inches(0.2)
content_w = col_w*2 + gap_h - arrow_w - Inches(0.3)
add_text_box(s, content_left, card_bot_top + Inches(0.15), content_w, Inches(0.25),
             "DELIBERATION", font=SANS, size=9, color=CORAL, bold=True, uppercase=True)
add_text_box(s, content_left, card_bot_top + Inches(0.4), Inches(0.55), Inches(0.4),
             "03", font=SERIF, size=22, color=CORAL)
add_text_box(s, content_left + Inches(0.55), card_bot_top + Inches(0.4), content_w - Inches(0.55), Inches(0.28),
             "Facilitating Discussion", font=SANS, size=14, color=CREAM, bold=True)
add_text_box(s, content_left + Inches(0.55), card_bot_top + Inches(0.7), content_w - Inches(0.55), Inches(0.4),
             "Block captains facilitate time-bound community deliberation at the block level.",
             font=SANS, size=10, color=CREAM, spacing=1.35)
add_text_box(s, content_left + Inches(0.55), card_bot_top + card_h_bot - Inches(0.3), content_w - Inches(0.55), Inches(0.2),
             "WHO  Block captains lead · Altadena public participates · local facilitators support",
             font=SANS, size=8, color=RGBColor(0xC8, 0xCC, 0xC0), italic=True)

add_footer(s)

# ----- SLIDE 6: FAIR OAKS (hand-hacked) -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "Worked example")
add_slide_title(s, "Fair Oaks Corridor (Immediate need)", size=28)

# Same layout as slide 5 with Fair Oaks-specific content
down_steps = [
    ("01", "Definition",
     "CIP surveys customized for Altadena — by area, across six categories (Mobility · Stormwater · Sewer · Water · Energy · Multi-Benefit). Plain-language briefs in EN + ES.",
     None),
    ("02", "Training & Distribution",
     "Altagether's 150+ captains + commercial-owner network + QR signage at storefronts and yards.",
     None),
]
up_steps = [
    ("05", "Shared record",
     "Structured packets to CIP — manually plotted on map. No automated dashboard yet.",
     None),
    ("04", "Aggregation",
     "Decision-record cards · shared Google docs · per-category compilation. Manual — no dashboard yet.",
     None),
]
card_top = Inches(1.7)
card_h_top = Inches(3.2)
gap_v = Inches(0.12)
card_bot_top = card_top + card_h_top + gap_v
card_h_bot = Inches(1.3)
draw_flow_card(s, col_left, card_top, col_w, card_h_top, CORAL_TINT, CORAL, "↓",
                "↓ Content flows down", down_steps)
draw_flow_card(s, col_right, card_top, col_w, card_h_top, SAGE_TINT, SAGE, "↑",
                "↑ Consensus flows up", up_steps)
# Deliberation bottom
add_rounded_rect(s, col_left, card_bot_top, col_w*2 + gap_h, card_h_bot, FOREST)
add_rounded_rect(s, col_left, card_bot_top, arrow_w, card_h_bot, FOREST_SOFT)
for di in range(2):
    add_text_box(s, col_left, card_bot_top + Inches(0.3) + di*Inches(0.3), arrow_w, Inches(0.25),
                 "•", font=SANS, size=18, color=CREAM, bold=True, align="center")
content_left = col_left + arrow_w + Inches(0.2)
content_w = col_w*2 + gap_h - arrow_w - Inches(0.3)
add_text_box(s, content_left, card_bot_top + Inches(0.12), content_w, Inches(0.22),
             "DELIBERATION", font=SANS, size=9, color=CORAL, bold=True, uppercase=True)
add_text_box(s, content_left, card_bot_top + Inches(0.35), Inches(0.55), Inches(0.4),
             "03", font=SERIF, size=20, color=CORAL)
add_text_box(s, content_left + Inches(0.55), card_bot_top + Inches(0.35), content_w - Inches(0.55), Inches(0.28),
             "Facilitating Discussion", font=SANS, size=13, color=CREAM, bold=True)
add_text_box(s, content_left + Inches(0.55), card_bot_top + Inches(0.65), content_w - Inches(0.55), Inches(0.5),
             "Walking tours · Zoom calls · facilitated input sessions · Design Fair for commercial owners and adjacent residents.",
             font=SANS, size=10, color=CREAM, spacing=1.35)
# Footnote
fnote_top = card_bot_top + card_h_bot + Inches(0.1)
add_rounded_rect(s, col_left, fnote_top, col_w*2 + gap_h, Inches(0.4), SAGE_TINT)
add_text_box(s, col_left + Inches(0.2), fnote_top, col_w*2 + gap_h - Inches(0.4), Inches(0.4),
             "Fair Oaks is the hand-hacked version — the framework works while we build the automated tools (real-time dashboard, auto-aggregation) once funded.",
             font=SANS, size=10, color=FOREST, italic=True, align="center", vertical="middle")

add_footer(s)

# ----- SLIDE 7: PROJECT PLAN GANTT -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "Project plan")
add_slide_title(s, "What happens after endorsement", size=28)

# Same Gantt structure as slide 3
chart_top = Inches(1.8)
months_area_left = chart_left + label_col_width

# Year headers
year_h = Inches(0.25)
add_rect(s, months_area_left, chart_top, year_2026_width - Inches(0.02), year_h, SAGE_TINT)
add_text_box(s, months_area_left, chart_top, year_2026_width, year_h, "2026",
             font=SANS, size=9, color=FOREST, bold=True, align="center", vertical="middle", uppercase=True)
add_rect(s, months_area_left + year_2026_width, chart_top, year_2027_width, year_h, SAGE_TINT)
add_text_box(s, months_area_left + year_2026_width, chart_top, year_2027_width, year_h, "2027",
             font=SANS, size=9, color=FOREST, bold=True, align="center", vertical="middle", uppercase=True)

month_y = chart_top + year_h + Inches(0.05)
for i, m in enumerate(months):
    add_text_box(s, months_area_left + i*month_width, month_y, month_width, Inches(0.22), m,
                 font=SANS, size=8, color=FOREST_SOFT, bold=True, align="center", uppercase=True)
under_y = month_y + Inches(0.24)
under = s.shapes.add_connector(1, months_area_left, under_y, months_area_left + months_area_width, under_y)
under.line.color.rgb = SAGE_TINT
under.line.width = Pt(0.5)

rows_top = under_y + Inches(0.1)
row_h = Inches(0.32)
row_gap = Inches(0.06)

project_rows = [
    ("CAPS presentation", "Today", [(3, 4, "May 21", CORAL, "solid")]),
    ("ATC + Altagether ratification", "", [(4, 6, "Jun–Jul", CORAL, "solid")]),
    ("EIFD advocacy", "At PFA hearings", [(5, 11, "Funding push", CORAL, "solid")]),
    ("Framework + toolkit build", "Standard, templates", [(4, 9, "Build", SAGE, "solid")]),
    ("Captain training", "Altagether facilitation", [(6, 10, "Training cohorts", SAGE, "solid")]),
    ("Fair Oaks pilot", "Hand-hacked live use case", [(7, 17, "Pilot in flight", FOREST, "solid")]),
    ("Popup retail", "Fair Oaks corridor activation", [(6, 23, "Activations & community-led commerce", PINK, "solid")]),
    ("Dashboard MVP", "Public accountability layer", [(13, 19, "Build & launch", SAGE, "solid")]),
    ("Iterate from pilot", "", [(14, 18, "Refine", SAGE, "solid")]),
    ("Scale to additional corridors", "", [(17, 23, "Roll out", CORAL, "solid")]),
]

for ri, (label, sublabel, bars) in enumerate(project_rows):
    y = rows_top + ri*(row_h + row_gap)
    add_text_box(s, chart_left, y, label_col_width - Inches(0.1), Inches(0.18), label,
                 font=SANS, size=10, color=FOREST, bold=True, vertical="middle")
    if sublabel:
        add_text_box(s, chart_left, y + Inches(0.16), label_col_width - Inches(0.1), Inches(0.15), sublabel,
                     font=SANS, size=7, color=FOREST_SOFT, vertical="middle")
    add_rounded_rect(s, months_area_left, y, months_area_width, row_h, RGBColor(0xEE, 0xF0, 0xEA))
    for (col_start, col_end, text, color, style) in bars:
        bar_left = months_area_left + (col_start - 1) * month_width + Inches(0.01)
        bar_right = months_area_left + (col_end - 1) * month_width - Inches(0.01)
        bar_w = bar_right - bar_left
        bar_text_color = WHITE
        if style == "solid":
            add_rounded_rect(s, bar_left, y + Inches(0.02), bar_w, row_h - Inches(0.04), color)
            add_text_box(s, bar_left + Inches(0.06), y + Inches(0.02), bar_w - Inches(0.12), row_h - Inches(0.04),
                         text, font=SANS, size=8, color=WHITE, bold=True, vertical="middle")
add_footer(s)

# ----- SLIDE 8: WHO PAYS -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "Who pays")
add_slide_title(s, "Who pays for what", size=32)
who_pays = [
    ("Sponsors", "Pay for their own delivery — out of existing project budgets, not new spend."),
    ("EIFD", "Funds the engagement infrastructure — administered through the PFA, a small fraction of the $2B financing layer."),
    ("Local professionals", "Do the work — facilitators, mediators, legal, translators from inside Altadena, so dollars recirculate."),
]
row_top = Inches(2.2)
row_h = Inches(1.3)
for i, (name, desc) in enumerate(who_pays):
    y = row_top + i*row_h
    add_text_box(s, Inches(0.7), y, Inches(4.5), Inches(0.45), name,
                 font=SANS, size=22, color=FOREST, bold=True)
    add_text_box(s, Inches(0.7), y + Inches(0.5), Inches(11.93), Inches(0.65), desc,
                 font=SANS, size=14, color=FOREST, spacing=1.4)
    if i < len(who_pays) - 1:
        ln = s.shapes.add_connector(1, Inches(0.7), y + row_h - Inches(0.05), Inches(12.63), y + row_h - Inches(0.05))
        ln.line.color.rgb = SAGE_TINT
        ln.line.width = Pt(0.5)
add_footer(s)

# ----- SLIDE 9: TEAM & PARTNERS -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, CREAM)
add_slide_tag(s, "Team & partners")
add_slide_title(s, "Team & partners", size=32)

# Collaboration team
add_text_box(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.25),
             "COLLABORATION TEAM", font=SANS, size=10, color=CORAL, bold=True, uppercase=True)
ln = s.shapes.add_connector(1, Inches(0.7), Inches(2.15), Inches(12.63), Inches(2.15))
ln.line.color.rgb = CORAL
ln.line.width = Pt(1.5)

team_orgs = [
    ("Altadena Village", "Esther Kim", CORAL),
    ("CoLab", "Megan Hinchliffe Gerig · Gabriella Caparco", SAGE),
    ("Alting", "Petra Wennberg", CORAL),
    ("Altagether", "John Mayo · 150+ block captains", FOREST),
]
card_w = Inches(2.95)
card_gap = Inches(0.15)
card_top = Inches(2.3)
for i, (org, lead, accent) in enumerate(team_orgs):
    x = Inches(0.7) + i*(card_w + card_gap)
    add_rounded_rect(s, x, card_top, card_w, Inches(1.2), WHITE)
    add_rect(s, x, card_top, card_w, Inches(0.08), accent)
    add_text_box(s, x + Inches(0.2), card_top + Inches(0.2), card_w - Inches(0.4), Inches(0.4),
                 org, font=SANS, size=16, color=FOREST, bold=True)
    add_text_box(s, x + Inches(0.2), card_top + Inches(0.6), card_w - Inches(0.4), Inches(0.55),
                 lead, font=SANS, size=11, color=FOREST_SOFT, spacing=1.4)

# Sponsors
add_text_box(s, Inches(0.7), Inches(3.85), Inches(12), Inches(0.25),
             "SPONSORS", font=SANS, size=10, color=SAGE, bold=True, uppercase=True)
ln = s.shapes.add_connector(1, Inches(0.7), Inches(4.15), Inches(12.63), Inches(4.15))
ln.line.color.rgb = SAGE
ln.line.width = Pt(1.5)

sponsors = [
    ("Altadena Bloom", "Ty Garrison"),
    ("Altadena Chamber of Commerce", "Judy Matthews"),
    ("Altadena Collective", "Chris Corbett"),
    ("Altadena Green", "Stephanie Landregan · Wynn Wilson"),
    ("Altadena Heritage", "Hans Allhoff · Michele Zack"),
    ("Altadena Town Council", "Dorothy Wong · Morgan Whirledge"),
    ("ARRC", "Anders Corey"),
    ("CCAR", "Hans Allhoff"),
    ("Dena Hive", "Dorothy Wong"),
    ("League of Women Voters", "Joan Riback"),
    ("Neighbors Building a Greater Altadena", "Karen Gibson"),
    ("NOMA", "Steve Lewis"),
    ("PUSD School Board", "Jennifer Hall Lee"),
    ("Steadfast", "Jodi McLaughlin · Hans Allhoff"),
    ("WRTT", "Ty Garrison"),
]
# 5 columns × 3 rows
chip_w = Inches(2.45)
chip_h = Inches(0.6)
chip_gap = Inches(0.05)
chip_top = Inches(4.3)
for i, (name, person) in enumerate(sponsors):
    row = i // 5
    col = i % 5
    x = Inches(0.7) + col*(chip_w + chip_gap)
    y = chip_top + row*(chip_h + chip_gap)
    add_rounded_rect(s, x, y, chip_w, chip_h, SAGE_TINT)
    add_text_box(s, x + Inches(0.12), y + Inches(0.08), chip_w - Inches(0.24), Inches(0.22),
                 name, font=SANS, size=9, color=FOREST, bold=True)
    add_text_box(s, x + Inches(0.12), y + Inches(0.3), chip_w - Inches(0.24), Inches(0.22),
                 person, font=SANS, size=8, color=FOREST_SOFT)

add_footer(s)

# ----- SLIDE 10: ASK -----
s = prs.slides.add_slide(blank_layout)
set_slide_bg(s, FOREST)
add_text_box(s, Inches(0.7), Inches(0.7), Inches(12), Inches(0.3),
             "WHAT WE WANT FROM YOU TODAY", font=SANS, size=11, color=CORAL, bold=True, uppercase=True)
add_text_box(s, Inches(0.7), Inches(1.7), Inches(11.93), Inches(2.5),
             "Endorsement of the approach — co-signing this process.",
             font=SERIF, size=44, color=CREAM, spacing=1.05)
add_text_box(s, Inches(0.7), Inches(4.3), Inches(11.5), Inches(1.4),
             "Not a separate letter to County. Not a budget commitment. CAPS standing alongside Altadena Village, CoLab, Alting, and Altagether in saying this is how community engagement should work. That weight helps move EIFD toward an earmark for the engagement infrastructure.",
             font=SANS, size=14, color=RGBColor(0xD8, 0xD2, 0xC8), spacing=1.5)
# Divider
ln = s.shapes.add_connector(1, Inches(0.7), Inches(5.85), Inches(12.63), Inches(5.85))
ln.line.color.rgb = RGBColor(0x4A, 0x5A, 0x42)
ln.line.width = Pt(0.5)
add_text_box(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.25),
             "WHAT ENDORSEMENT ENABLES NEXT", font=SANS, size=10, color=CORAL, bold=True, uppercase=True)
add_text_box(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.6),
             "Joint visibility at PFA hearings  ·  your co-signature on the Fair Oaks pilot framing  ·  a credible voice when this lands at ATC",
             font=SANS, size=12, color=RGBColor(0xD8, 0xD2, 0xC8), spacing=1.4)

# ============ SAVE ============
output = "/Users/PetraWCesario/Freya/alting/altadena-transparency/strategy/caps-2026-05-21/deck.pptx"
prs.save(output)
print(f"Saved: {output}")
